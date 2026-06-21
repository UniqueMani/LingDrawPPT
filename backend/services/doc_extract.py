from __future__ import annotations

import shutil
import subprocess
import tempfile
import time
from io import BytesIO
from pathlib import Path
from textwrap import wrap
from typing import Dict, List
from uuid import uuid4

from pypdf import PdfReader
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

from backend.config import PREVIEW_DIR


PREVIEW_ROOT = Path(PREVIEW_DIR)
PPTX_TO_PDF_TIMEOUT_SECONDS = 90
PP_SAVE_AS_PDF = 32


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for candidate in candidates:
        try:
            if Path(candidate).exists():
                return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _save_preview_pair(img: Image.Image, out_dir: Path, page: int) -> Dict[str, str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    preview_name = f"page-{page:03d}.png"
    thumb_name = f"thumb-{page:03d}.png"
    preview_path = out_dir / preview_name
    thumb_path = out_dir / thumb_name

    rgb = img.convert("RGB")
    rgb.save(preview_path, "PNG", optimize=True)

    thumb = rgb.copy()
    thumb.thumbnail((320, 200))
    canvas = Image.new("RGB", (320, 200), "#f3f1f2")
    x = (320 - thumb.width) // 2
    y = (200 - thumb.height) // 2
    canvas.paste(thumb, (x, y))
    canvas.save(thumb_path, "PNG", optimize=True)

    rel = out_dir.relative_to(PREVIEW_ROOT).as_posix()
    return {
        "preview_url": f"/previews/{rel}/{preview_name}",
        "thumbnail_url": f"/previews/{rel}/{thumb_name}",
        "preview_updated_at": int(time.time() * 1000),
    }


def _text_preview_image(title: str, text: str, *, wide: bool = False) -> Image.Image:
    width, height = (1280, 720) if wide else (900, 1200)
    img = Image.new("RGB", (width, height), "#ffffff")
    draw = ImageDraw.Draw(img)
    title_font = _font(42 if wide else 36, bold=True)
    body_font = _font(24 if wide else 22)
    meta_font = _font(18)

    accent = "#8b2942"
    draw.rectangle((0, 0, width, 18), fill=accent)
    draw.text((56, 54), title or "未命名页面", fill="#1a1a1a", font=title_font)
    draw.text((56, 110), "文本降级预览", fill=accent, font=meta_font)

    x, y = 56, 170
    max_chars = 42 if wide else 34
    line_h = 38 if wide else 34
    max_lines = (height - y - 56) // line_h
    lines: List[str] = []
    for raw in (text or "").splitlines():
        raw = raw.strip()
        if not raw:
            continue
        lines.extend(wrap(raw, width=max_chars) or [raw])
        if len(lines) >= max_lines:
            break
    for line in lines[:max_lines]:
        draw.text((x, y), line, fill="#2d2d2d", font=body_font)
        y += line_h
    return img


def _render_pdf_previews(file_bytes: bytes, pages_detail: List[Dict[str, object]], out_dir: Path) -> None:
    try:
        import fitz  # type: ignore
    except Exception:
        for detail in pages_detail:
            page = int(detail.get("page", 1) or 1)
            urls = _save_preview_pair(
                _text_preview_image(str(detail.get("title", "")), str(detail.get("text", ""))),
                out_dir,
                page,
            )
            detail.update(urls)
        return

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for idx, detail in enumerate(pages_detail):
        page_no = int(detail.get("page", idx + 1) or (idx + 1))
        try:
            page = doc.load_page(page_no - 1)
            pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        except Exception:
            img = _text_preview_image(str(detail.get("title", "")), str(detail.get("text", "")), wide=True)
        detail.update(_save_preview_pair(img, out_dir, page_no))
    doc.close()


def _convert_pptx_to_pdf_with_powerpoint(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except Exception:
        return False

    app = None
    presentation = None
    pythoncom.CoInitialize()
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = app.Presentations.Open(str(pptx_path), False, False, False)
        presentation.SaveAs(str(pdf_path), PP_SAVE_AS_PDF)
        return pdf_path.exists() and pdf_path.stat().st_size > 0
    except Exception:
        return False
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def _soffice_path() -> str | None:
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    common = Path("C:/Program Files/LibreOffice/program/soffice.exe")
    if common.exists():
        return str(common)
    return None


def _convert_pptx_to_pdf_with_libreoffice(pptx_path: Path, pdf_path: Path) -> bool:
    soffice = _soffice_path()
    if not soffice:
        return False
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_path.parent),
            str(pptx_path),
        ],
        capture_output=True,
        text=True,
        timeout=PPTX_TO_PDF_TIMEOUT_SECONDS,
        check=False,
    )
    converted = pdf_path.parent / f"{pptx_path.stem}.pdf"
    if converted.exists() and converted != pdf_path:
        converted.replace(pdf_path)
    return result.returncode == 0 and pdf_path.exists() and pdf_path.stat().st_size > 0


def _convert_pptx_to_pdf(file_bytes: bytes) -> bytes | None:
    with tempfile.TemporaryDirectory(prefix="lingdraw_pptx_") as tmp:
        tmp_dir = Path(tmp)
        pptx_path = tmp_dir / "source.pptx"
        pdf_path = tmp_dir / "source.pdf"
        pptx_path.write_bytes(file_bytes)
        if _convert_pptx_to_pdf_with_powerpoint(pptx_path, pdf_path):
            return pdf_path.read_bytes()
        if _convert_pptx_to_pdf_with_libreoffice(pptx_path, pdf_path):
            return pdf_path.read_bytes()
    return None


def _render_pptx_previews(file_bytes: bytes, pages_detail: List[Dict[str, object]], out_dir: Path) -> None:
    pdf_bytes = _convert_pptx_to_pdf(file_bytes)
    if pdf_bytes:
        _render_pdf_previews(pdf_bytes, pages_detail, out_dir)
        return

    for detail in pages_detail:
        page = int(detail.get("page", 1) or 1)
        img = _text_preview_image(str(detail.get("title", "")), str(detail.get("text", "")), wide=True)
        detail.update(_save_preview_pair(img, out_dir, page))


def _attach_previews(filename: str, file_bytes: bytes, pages_detail: List[Dict[str, object]], out_dir: Path | None = None) -> None:
    if not pages_detail:
        return
    suffix = Path(filename).suffix.lower()
    target_dir = out_dir or (PREVIEW_ROOT / uuid4().hex)
    target_dir.mkdir(parents=True, exist_ok=True)
    if suffix == ".pdf":
        _render_pdf_previews(file_bytes, pages_detail, target_dir)
    elif suffix == ".pptx":
        _render_pptx_previews(file_bytes, pages_detail, target_dir)


def refresh_previews_in_dir(filename: str, file_bytes: bytes, pages_detail: List[Dict[str, object]], out_dir: Path) -> None:
    if not pages_detail:
        return
    out_dir.mkdir(parents=True, exist_ok=True)
    _attach_previews(filename, file_bytes, pages_detail, out_dir=out_dir)


def _page_detail_item(pages_detail: List[Dict[str, object]], page: int) -> Dict[str, object] | None:
    for item in pages_detail:
        if int(item.get("page", 0) or 0) == page:
            return item
    return None


def _render_single_pdf_page_preview(
    pdf_bytes: bytes,
    detail: Dict[str, object],
    out_dir: Path,
    page: int,
) -> None:
    page_no = int(detail.get("page", page) or page)
    try:
        import fitz  # type: ignore
    except Exception:
        img = _text_preview_image(str(detail.get("title", "")), str(detail.get("text", "")), wide=True)
        detail.update(_save_preview_pair(img, out_dir, page_no))
        return

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        try:
            page_obj = doc.load_page(page_no - 1)
            pix = page_obj.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        except Exception:
            img = _text_preview_image(str(detail.get("title", "")), str(detail.get("text", "")), wide=True)
        detail.update(_save_preview_pair(img, out_dir, page_no))
    finally:
        doc.close()


def refresh_single_page_preview_in_dir(
    filename: str,
    file_bytes: bytes,
    pages_detail: List[Dict[str, object]],
    out_dir: Path,
    page: int,
) -> None:
    if not pages_detail:
        return
    detail = _page_detail_item(pages_detail, page)
    if detail is None:
        return

    out_dir.mkdir(parents=True, exist_ok=True)
    suffix = Path(filename).suffix.lower()
    if suffix == ".pptx":
        pdf_bytes = _convert_pptx_to_pdf(file_bytes)
        if pdf_bytes:
            _render_single_pdf_page_preview(pdf_bytes, detail, out_dir, page)
            return
    if suffix == ".pdf":
        _render_single_pdf_page_preview(file_bytes, detail, out_dir, page)
        return

    img = _text_preview_image(str(detail.get("title", "")), str(detail.get("text", "")), wide=True)
    detail.update(_save_preview_pair(img, out_dir, page))


def export_document_bytes(
    *,
    original_filename: str,
    file_bytes: bytes,
    export_format: str,
) -> tuple[bytes, str, str]:
    suffix = Path(original_filename).suffix.lower()
    stem = Path(original_filename).stem or "document"
    fmt = export_format.lower()

    if fmt in {"pptx", "ppt"}:
        if suffix != ".pptx":
            raise ValueError("当前文件不是 PPTX，无法导出为 PPT")
        return (
            file_bytes,
            original_filename if original_filename.lower().endswith(".pptx") else f"{stem}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    if fmt == "pdf":
        if suffix == ".pdf":
            return file_bytes, original_filename, "application/pdf"
        if suffix != ".pptx":
            raise ValueError("当前文件不支持导出为 PDF")
        pdf_bytes = _convert_pptx_to_pdf(file_bytes)
        if not pdf_bytes:
            raise RuntimeError("PDF 转换失败，请确认已安装 LibreOffice 或 Microsoft PowerPoint")
        return pdf_bytes, f"{stem}.pdf", "application/pdf"

    raise ValueError("不支持的导出格式")


def _normalize_lines(lines: List[str], max_lines: int = 240) -> str:
    cleaned: List[str] = []
    for line in lines:
        s = (line or "").strip()
        if not s:
            continue
        cleaned.append(s)
        if len(cleaned) >= max_lines:
            break
    return "\n".join(cleaned)


def _normalized_text_block(
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    page_width: float,
    page_height: float,
) -> Dict[str, object] | None:
    cleaned = (text or "").strip()
    if not cleaned or page_width <= 0 or page_height <= 0:
        return None
    return {
        "text": cleaned,
        "x": max(0.0, min(1.0, x / page_width)),
        "y": max(0.0, min(1.0, y / page_height)),
        "width": max(0.0, min(1.0, width / page_width)),
        "height": max(0.0, min(1.0, height / page_height)),
    }


def _extract_pdf_text_blocks(file_bytes: bytes) -> List[List[Dict[str, object]]]:
    try:
        import fitz  # type: ignore
    except Exception:
        return []

    result: List[List[Dict[str, object]]] = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        for page in doc:
            blocks: List[Dict[str, object]] = []
            for x0, y0, x1, y1, text, *_ in page.get_text("blocks"):
                block = _normalized_text_block(
                    str(text),
                    float(x0),
                    float(y0),
                    float(x1 - x0),
                    float(y1 - y0),
                    float(page.rect.width),
                    float(page.rect.height),
                )
                if block:
                    blocks.append(block)
            result.append(blocks)
    finally:
        doc.close()
    return result


def _shape_text(shape: object) -> str:
    if getattr(shape, "has_table", False):
        rows: List[str] = []
        for row in shape.table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if line:
                rows.append(line)
        return "\n".join(rows)
    return str(getattr(shape, "text", "") or "").strip()


def extract_text_from_pdf(file_bytes: bytes) -> Dict[str, object]:
    reader = PdfReader(BytesIO(file_bytes))
    blocks_by_page = _extract_pdf_text_blocks(file_bytes)
    lines: List[str] = []
    pages_detail: List[Dict[str, object]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            lines.append(text)
        title = (text.splitlines()[0].strip() if text.splitlines() else f"第 {idx} 页")[:80]
        pages_detail.append(
            {
                "page": idx,
                "title": title,
                "text": text.strip(),
                "text_blocks": blocks_by_page[idx - 1] if idx <= len(blocks_by_page) else [],
            }
        )
    text = _normalize_lines(lines)
    first_line = text.splitlines()[0].strip() if text else ""
    return {
        "text": text,
        "title": first_line[:80],
        "pages": len(reader.pages),
        "pages_detail": pages_detail,
    }


def extract_text_from_pptx(file_bytes: bytes) -> Dict[str, object]:
    prs = Presentation(BytesIO(file_bytes))
    page_width = float(prs.slide_width)
    page_height = float(prs.slide_height)
    lines: List[str] = []
    pages_detail: List[Dict[str, object]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: List[str] = []
        text_blocks: List[Dict[str, object]] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text:
                lines.append(text)
                slide_lines.append(text)
                block = _normalized_text_block(
                    text,
                    float(shape.left),
                    float(shape.top),
                    float(shape.width),
                    float(shape.height),
                    page_width,
                    page_height,
                )
                if block:
                    text_blocks.append(block)
        joined = "\n".join(slide_lines)
        title = (slide_lines[0] if slide_lines else f"第 {idx} 页")[:80]
        pages_detail.append(
            {
                "page": idx,
                "title": title,
                "text": joined,
                "text_blocks": text_blocks,
                "page_width": page_width,
                "page_height": page_height,
            }
        )
    text = _normalize_lines(lines)
    first_line = text.splitlines()[0].strip() if text else ""
    return {
        "text": text,
        "title": first_line[:80],
        "pages": len(prs.slides),
        "pages_detail": pages_detail,
    }


def extract_text_from_document(filename: str, file_bytes: bytes) -> Dict[str, object]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        result = extract_text_from_pdf(file_bytes)
        _attach_previews(filename, file_bytes, list(result.get("pages_detail", []) or []))
        return result
    if suffix == ".pptx":
        result = extract_text_from_pptx(file_bytes)
        _attach_previews(filename, file_bytes, list(result.get("pages_detail", []) or []))
        return result
    raise ValueError("仅支持 .pdf 或 .pptx 文件")
