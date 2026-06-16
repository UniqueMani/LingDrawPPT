from __future__ import annotations

import base64
import json
import logging
import time
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Tuple
from uuid import uuid4

import httpx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.config import PREVIEW_DIR, UPLOAD_DIR
from backend.services.doc_extract import refresh_previews_in_dir, refresh_single_page_preview_in_dir
from backend.services.ppt_placement import compute_recommended_placement

logger = logging.getLogger(__name__)

STAGED_DIR = Path(UPLOAD_DIR) / "staged"


def is_pptx_filename(name: str) -> bool:
    return Path(name or "").suffix.lower() == ".pptx"


def _decode_data_url(image_data: str) -> bytes:
    raw = (image_data or "").strip()
    if not raw:
        raise ValueError("图片数据为空")
    if raw.startswith("data:"):
        if "," not in raw:
            raise ValueError("无效的图片 data URL")
        raw = raw.split(",", 1)[1]
    content = base64.b64decode(raw, validate=False)
    if not content:
        raise ValueError("图片内容为空")
    return content


def save_staged_image(user_id: int, image_data: str) -> str:
    STAGED_DIR.mkdir(parents=True, exist_ok=True)
    content = _decode_data_url(image_data)
    token = f"{user_id}_{uuid4().hex}.png"
    path = STAGED_DIR / token
    path.write_bytes(content)
    return f"/api/ppt/staged/{token}"


def read_staged_image(name: str) -> bytes:
    safe = Path(name).name
    path = (STAGED_DIR / safe).resolve()
    root = STAGED_DIR.resolve()
    if not str(path).startswith(str(root)):
        raise ValueError("无效的图片路径")
    if not path.is_file():
        raise ValueError("暂存图片不存在")
    return path.read_bytes()


async def download_image_bytes(image_url: str) -> bytes:
    url = (image_url or "").strip()
    if not url:
        raise ValueError("image_url 不能为空")
    if url.startswith("data:"):
        return _decode_data_url(url)
    if "/api/ppt/staged/" in url:
        name = url.rstrip("/").split("/")[-1]
        return read_staged_image(name)
    async with httpx.AsyncClient(timeout=120.0, follow_redirects=True) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        content = resp.content
    if not content:
        raise ValueError("图片下载失败：内容为空")
    return content

def load_pages_detail(raw: str) -> List[Dict[str, object]]:
    try:
        data = json.loads(raw or "[]")
    except json.JSONDecodeError:
        return []
    return list(data) if isinstance(data, list) else []


def page_detail(pages_detail: List[Dict[str, object]], page: int) -> Dict[str, object]:
    for item in pages_detail or []:
        if int(item.get("page", 0) or 0) == page:
            return item
    raise ValueError(f"未找到第 {page} 页")


def preview_dir_from_pages_detail(pages_detail: List[Dict[str, object]]) -> Path | None:
    preview_root = Path(PREVIEW_DIR)
    for item in pages_detail:
        preview_url = str(item.get("preview_url") or "")
        parts = preview_url.strip("/").split("/")
        if len(parts) >= 2 and parts[0] == "previews":
            return preview_root / parts[1]
    return None


def recommend_placement_for_page(
    pages_detail: List[Dict[str, object]],
    *,
    page: int,
    aspect_ratio: str,
) -> Tuple[Dict[str, float], List[Dict[str, object]], Dict[str, object]]:
    detail = page_detail(pages_detail, page)
    text_blocks = list(detail.get("text_blocks") or [])
    page_width = float(detail.get("page_width") or 0)
    page_height = float(detail.get("page_height") or 0)
    slide_aspect = page_width / page_height if page_width > 0 and page_height > 0 else 16.0 / 9.0
    recommended, occupied = compute_recommended_placement(
        text_blocks,
        aspect_ratio=aspect_ratio,
        slide_aspect=slide_aspect,
    )
    return recommended, occupied, detail


def normalized_rect_to_emu(
    placement: Dict[str, float],
    slide_width: int,
    slide_height: int,
) -> Tuple[int, int, int, int]:
    left = int(float(placement["x"]) * slide_width)
    top = int(float(placement["y"]) * slide_height)
    width = int(float(placement["width"]) * slide_width)
    height = int(float(placement["height"]) * slide_height)
    return left, top, width, height


def infer_aspect_ratio_label(pixel_aspect: float) -> str:
    if pixel_aspect <= 0:
        return "16:9"
    presets: List[Tuple[str, float]] = [
        ("16:9", 16.0 / 9.0),
        ("4:3", 4.0 / 3.0),
        ("1:1", 1.0),
        ("9:16", 9.0 / 16.0),
        ("21:9", 21.0 / 9.0),
    ]
    best = presets[0][0]
    best_diff = float("inf")
    for label, value in presets:
        diff = abs(pixel_aspect - value)
        if diff < best_diff:
            best_diff = diff
            best = label
    return best


def _shape_to_picture_info(shape, shape_index: int, slide_width: int, slide_height: int) -> Dict[str, object]:
    slide_aspect = slide_width / max(slide_height, 1)
    width = float(shape.width) / slide_width
    height = float(shape.height) / slide_height
    x = float(shape.left) / slide_width
    y = float(shape.top) / slide_height
    pixel_aspect = (width / max(height, 1e-6)) * slide_aspect
    return {
        "shape_index": shape_index,
        "x": max(0.0, min(1.0, x)),
        "y": max(0.0, min(1.0, y)),
        "width": max(0.0, min(1.0, width)),
        "height": max(0.0, min(1.0, height)),
        "aspect_ratio": infer_aspect_ratio_label(pixel_aspect),
    }


def list_page_pictures(*, pptx_path: Path, page: int) -> Tuple[List[Dict[str, object]], int, int]:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")
    prs = Presentation(str(pptx_path))
    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    pictures: List[Dict[str, object]] = []
    for idx, shape in enumerate(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        pictures.append(_shape_to_picture_info(shape, idx, slide_width, slide_height))
    return pictures, slide_width, slide_height


def read_picture_bytes(*, pptx_path: Path, page: int, shape_index: int) -> Tuple[bytes, str]:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")
    prs = Presentation(str(pptx_path))
    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError("无效的图片索引")
    shape = slide.shapes[shape_index]
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise ValueError("目标不是图片")

    content = shape.image.blob
    if not content:
        raise ValueError("图片内容为空")
    ext = str(getattr(shape.image, "ext", "") or "png").lower()
    media_type = "image/png" if ext == "png" else "image/jpeg"
    return content, media_type


def update_picture_placement(
    *,
    pptx_path: Path,
    page: int,
    shape_index: int,
    placement: Dict[str, float],
) -> None:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")
    prs = Presentation(str(pptx_path))
    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError("无效的图片索引")
    shape = slide.shapes[shape_index]
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise ValueError("目标不是图片")

    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    left, top, width, height = normalized_rect_to_emu(placement, slide_width, slide_height)
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    prs.save(str(pptx_path))


def remove_picture_by_index(*, pptx_path: Path, page: int, shape_index: int) -> bool:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")
    prs = Presentation(str(pptx_path))
    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return False
    shape = slide.shapes[shape_index]
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    element = shape._element
    element.getparent().remove(element)
    prs.save(str(pptx_path))
    return True


def remove_last_picture_from_pptx(*, pptx_path: Path, page: int) -> bool:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")
    prs = Presentation(str(pptx_path))
    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    for idx in range(len(slide.shapes) - 1, -1, -1):
        shape = slide.shapes[idx]
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element = shape._element
            element.getparent().remove(element)
            prs.save(str(pptx_path))
            return True
    return False


def insert_image_into_pptx(
    *,
    pptx_path: Path,
    page: int,
    image_bytes: bytes,
    placement: Dict[str, float],
) -> Dict[str, object]:
    if not pptx_path.is_file():
        raise ValueError("PPT 文件不存在")

    t0 = time.perf_counter()
    file_size = pptx_path.stat().st_size
    prs = Presentation(str(pptx_path))
    t_loaded = time.perf_counter()

    if page < 1 or page > len(prs.slides):
        raise ValueError(f"页码超出范围：{page}")

    slide = prs.slides[page - 1]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    left, top, width, height = normalized_rect_to_emu(
        placement,
        slide_width,
        slide_height,
    )
    picture = slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)
    t_added = time.perf_counter()

    prs.save(str(pptx_path))
    t_saved = time.perf_counter()

    logger.info(
        "insert_image_into_pptx page=%s file_size=%.1fKB load=%.2fs add=%.2fs save=%.2fs total=%.2fs",
        page,
        file_size / 1024,
        t_loaded - t0,
        t_added - t_loaded,
        t_saved - t_added,
        t_saved - t0,
    )

    shape_index = len(slide.shapes) - 1
    for idx, shape in enumerate(slide.shapes):
        if shape.shape_id == picture.shape_id:
            shape_index = idx
            break
    return _shape_to_picture_info(picture, shape_index, slide_width, slide_height)


def refresh_file_previews(
    *,
    original_filename: str,
    pptx_path: Path,
    pages_detail: List[Dict[str, object]],
) -> List[Dict[str, object]]:
    preview_dir = preview_dir_from_pages_detail(pages_detail)
    if preview_dir is None:
        return pages_detail
    try:
        refresh_previews_in_dir(original_filename, pptx_path.read_bytes(), pages_detail, preview_dir)
    except Exception as exc:
        logger.warning("refresh previews after insert failed: %s", exc)
    return pages_detail


def refresh_page_preview(
    *,
    original_filename: str,
    pptx_path: Path,
    pages_detail: List[Dict[str, object]],
    page: int,
) -> List[Dict[str, object]]:
    preview_dir = preview_dir_from_pages_detail(pages_detail)
    if preview_dir is None:
        return pages_detail
    try:
        refresh_single_page_preview_in_dir(
            original_filename,
            pptx_path.read_bytes(),
            pages_detail,
            preview_dir,
            page,
        )
    except Exception as exc:
        logger.warning("refresh page preview failed: %s", exc)
    return pages_detail
